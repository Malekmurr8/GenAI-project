import openai
import os
import json
from pptx import Presentation
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from dotenv import load_dotenv

# ✅ Load environment variables
load_dotenv()

app = FastAPI()

# ✅ Securely load OpenAI API key
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

@app.get("/test-api-key")
def test_api_key():
    """Endpoint to check if the API key is loaded correctly."""
    return {"embedded_key": OPENAI_API_KEY}

def generate_content(topic, country, openai_api_key, model="gpt-3.5-turbo"):
    """Generate content using OpenAI's GPT model."""
    openai.api_key = openai_api_key
    system_prompt = "You are a helpful AI that returns only valid JSON without additional commentary."
    user_prompt = f"""
    {{
      "Title 2": "Short main title referencing topic/country",
      "Rectangle 25": "One-line heading for first sub-topic",
      "Rectangle 29": "One-line heading for second sub-topic",
      "Rectangle 35": "One-line heading for third sub-topic",
      "Rectangle 39": "Provide 2 sentences about first sub-topic. Make sure to include facts, numbers, and examples as much as possible",
      "Rectangle 40": "Provide 2 sentences about second sub-topic. Make sure to include facts, numbers, and examples as much as possible",
      "Rectangle 41": "Provide 2 sentences about third sub-topic. Make sure to include facts, numbers, and examples as much as possible"
    }}
    Instructions: Tailored to \"{topic}\" and \"{country}\".
    """
    response = openai.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.8,
        max_tokens=1400
    )
    content = response.choices[0].message.content
    if "```" in content:
        content = content.split("```")[1].strip()
    return json.loads(content)

FONT_SIZES = {
    "Title 2": 40, "Rectangle 25": 18, "Rectangle 29": 18, "Rectangle 35": 18,
    "Rectangle 39": 16, "Rectangle 40": 16, "Rectangle 41": 16
}
BLACK_SHAPES = {"Rectangle 39", "Rectangle 40", "Rectangle 41"}

def set_text_in_shape(shape, new_text, is_black, font_size_pt):
    """Update text inside a PowerPoint shape."""
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    color_rgb = RGBColor(0, 0, 0) if is_black else RGBColor(255, 255, 255)
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p.text = new_text
    p.font.size = Pt(font_size_pt)
    p.font.color.rgb = color_rgb
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

def update_texts(slide, new_data):
    """Update all text fields in the PowerPoint slide."""
    for shape in slide.shapes:
        name = shape.name
        if name in new_data and name in FONT_SIZES:
            is_black = name in BLACK_SHAPES
            font_size_pt = FONT_SIZES[name]
            set_text_in_shape(shape, new_data[name], is_black, font_size_pt)

def replace_flag(slide, country_name):
    """Replace flag image based on the selected country."""
    flags_dir = "flags"
    filename = f"{country_name.lower()}.png"
    flag_path = os.path.join(flags_dir, filename)
    if not os.path.exists(flag_path):
        print(f"⚠️ Flag file not found: {flag_path}")
        return
    for shape in slide.shapes:
        if shape.name == "Picture 5" and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(flag_path, left, top, width, height)
            return

class SlideRequest(BaseModel):
    topic: str
    country: str

@app.post("/generate-slide")
def generate_slide(request: SlideRequest):
    try:
        print(f"🔹 Received Request: Topic={request.topic}, Country={request.country}")

        # ✅ Ensure 'slides/' directory exists
        output_folder = "slides"
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # ✅ Generate AI-based slide content
        new_data = generate_content(request.topic, request.country, OPENAI_API_KEY)

        # ✅ Load the template
        template_path = "templates/Australia Benchmark.pptx"
        if not os.path.exists(template_path):
            raise HTTPException(status_code=500, detail=f"Template file not found: {template_path}")

        prs = Presentation(template_path)
        slide = prs.slides[0]

        # ✅ Apply AI-generated text updates
        update_texts(slide, new_data)

        # ✅ Replace the flag image
        replace_flag(slide, request.country)

        # ✅ Save the modified slide
        output_path = f"{output_folder}/{request.topic}_{request.country}.pptx"
        prs.save(output_path)

        print(f"✅ Slide saved: {output_path}")
        return FileResponse(output_path, filename=os.path.basename(output_path))

    except Exception as e:
        print(f"❌ Error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
